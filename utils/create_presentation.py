import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from bot.settings import PRESENTATION_PATH, PATH_TO_0_IMG


def create_slide(text: str) -> None:
    presentation_path = PRESENTATION_PATH

    """Открываем шаблон презентации"""
    prs = Presentation(presentation_path)
    prs.save('given.pptx')
    prs = Presentation('given.pptx')
    """Добавляем слайд с макетом"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    """Добавляем текст из сообщения вверху слайда"""
    first_shape = shapes.add_textbox(
        left=Inches(0.1),
        top=Inches(0.1),
        width=Inches(6),
        height=Inches(0.47)
    )
    first_shape_frame = first_shape.text_frame
    first_shape_frame.word_wrap = True
    p = first_shape_frame.paragraphs[0]
    p.text = text  # изменяемое значение (текст из сообщения)
    p.font.name = 'Newfont'
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(190, 190, 190)

    """Добавляем черный прямоугольник"""
    second_shape = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0),
        top=Inches(0.63),
        width=prs.slide_width,
        height=Inches(4.25)
    )
    fill_color = second_shape.fill
    fill_color.solid()
    fill_color.fore_color.rgb = RGBColor(0, 0, 0)
    line = second_shape.line
    line.color.rgb = RGBColor(0, 0, 0)

    """Добавляем нулевую картинку"""
    third_shape = shapes.add_picture(
        PATH_TO_0_IMG,
        left=Inches(1.45),
        top=Inches(4.645),
        width=Inches(8.45),
        height=Inches(0.24)
    )

    """Добавляем текст "Освещенность, lx"""
    fourth_shape = shapes.add_textbox(
        left=Inches(0.2),
        top=Inches(4.65),
        width=Inches(1.25),
        height=Inches(0.25)
    )
    fourth_shape_frame = fourth_shape.text_frame
    first_shape_frame.word_wrap = True
    p = fourth_shape_frame.paragraphs[0]
    p.text = "Освещенность, lx"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = RGBColor(190, 190, 190)

    """Добавляем порядковый номер слайда"""
    fifth_shape = shapes.add_textbox(
        left=Inches(0.2),
        top=Inches(4.65),
        width=Inches(0.5),
        height=Inches(0.75)
    )
    fifth_shape_frame = fifth_shape.text_frame
    p = fifth_shape_frame.paragraphs[0]
    p.text = "3"  # изменяемое значение (номер слайда)
    p.font.size = Pt(60)
    p.font.color.rgb = RGBColor(190, 190, 190)
    p.alignment = PP_ALIGN.CENTER

    """Добавляем 2 картинки"""
    images = os.listdir('/Users/aleksejdegtarev/PycharmProjects/PresentationBot/bot/images/temp')
    if len(images) == 2:
        shapes.add_picture(
            f'/Users/aleksejdegtarev/PycharmProjects/PresentationBot/bot/images/temp/{images[0]}',
            left=Inches(0),
            top=Inches(0.63),
            width=prs.slide_width / 2,
            height=Inches(4.02)
        )
        shapes.add_picture(
            f'/Users/aleksejdegtarev/PycharmProjects/PresentationBot/bot/images/temp/{images[1]}',
            left=prs.slide_width / 2,
            top=Inches(0.63),
            width=prs.slide_width / 2,
            height=Inches(4.02)
        )
    elif len(images) == 4:
        shapes.add_picture(
            f'/Users/aleksejdegtarev/PycharmProjects/PresentationBot/bot/images/temp/{images[0]}',
            left=Inches(0),
            top=Inches(0.75),
            width=Inches(3.5),
            height=Inches(2.63)
        )
        shapes.add_picture(
            f'/Users/aleksejdegtarev/PycharmProjects/PresentationBot/bot/images/temp/{images[1]}',
            left=Inches(3.6),
            top=Inches(0.75),
            width=Inches(3.5),
            height=Inches(2.63)
        )
        shapes.add_picture(
            f'/Users/aleksejdegtarev/PycharmProjects/PresentationBot/bot/images/temp/{images[2]}',
            left=Inches(7.2),
            top=Inches(0.75),
            width=Inches(2.46),
            height=Inches(1.85)
        )
        shapes.add_picture(
            f'/Users/aleksejdegtarev/PycharmProjects/PresentationBot/bot/images/temp/{images[3]}',
            left=Inches(7.2),
            top=Inches(2.7),
            width=Inches(2.46),
            height=Inches(1.85)
        )

    """Сохраняем презентацию"""
    prs.save('given.pptx')

    for item in images:
        os.remove(f'/Users/aleksejdegtarev/PycharmProjects/PresentationBot/bot/images/temp/{item}')
